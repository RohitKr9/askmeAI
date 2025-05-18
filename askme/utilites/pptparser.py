from pptx import Presentation

def getText(mypptx):
    ppt = Presentation(mypptx)
    content = ''
    #we will iterate through each slide to get text
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame
                for para in text.paragraphs:
                    for run in para.runs:
                        content += run.text

    print(content)
    return content

